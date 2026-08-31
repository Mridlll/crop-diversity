"""
84_build_deck.py

The crop diversity deck, CEEW house style, 13.333 by 7.5 inches.

Imports the house furniture rather than restating it. Figures come from
83_deck_figures.py and carry data only: the slide title names the page, the
strip carries the finding, the footnote carries definitions.

The running order is an argument the reader follows from ignorance to conclusion.
Act one asks what rides on a district's crop mix. Act two describes: the measure,
the national picture, the geography, the concentration in two cereals, and whether
the balance has moved. Act three opens on the irrigation hump, which is the pattern
that generates every candidate explanation, and screens thirteen candidates down to
five. Act four states those five. Act five tests them one at a time and keeps the
one that fails on the page with the four that held.

Page shapes vary deliberately. Eight of the fourteen pages carry a shape that
appears once: a lede over three panels with a statistics band, a summary-statistics
page of three blocks, a large map with a side rail, a full-width figure over a
statistics band, a hypothesis table, a figure over a table, a scorecard and a page
of two stacked tables.

VOCABULARY IS LOCKED. It is "the effective number of crops" everywhere, "crops
grown in an average year", "evenness", "the irrigation share of sown area",
"weekly haat", "mandi", "regular market" and "producer organisation". The composite
index appears once, in the footnote to slide two, as the reason a count is used.

Output: deck/crop_diversity.pptx
"""
import os
import sys
from pathlib import Path

from PIL import Image
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

FURN = r"D:/Alternative Proteins/demand_pathways/deck"
sys.path.insert(0, FURN)

from deck_furniture import (  # noqa: E402
    plain_slide, chrome, source, footnote, highlight_strip, titled_panel,
    panel, box, set_text, check_fits, wrapped_lines, stack_in, prs,
    RUNNING_ORDER, LINE_PITCH_IN, TWOUP_X, TWOUP_W, STRIP_X, STRIP_W,
    PANEL_ML, PANEL_MR, PANEL_MB, INK, MUTED, TINT, WHITE, ACCENT, LINE)

REPO = Path(r"D:/crop-diversity")
FIGS = REPO / "deck" / "figs"
OUT = REPO / "deck" / "crop_diversity.pptx"

SRC = "Ministry of Agriculture and Farmers Welfare; SHRUG 2.1, Development Data Lab"

# The three-across span, which is box_row's own left and right edges with its own
# gap, so a three-panel row lines up with the two-up pages and the four-slot band.
BAND_GAP = 0.34
THIRD_W = (12.88 - 0.43 - 2 * BAND_GAP) / 3          # 3.9233
THIRD_X = [0.43 + i * (THIRD_W + BAND_GAP) for i in range(3)]


def place(slide, name, x, y, max_w, max_h, centre=True):
    """Place a figure scaled to fit its box, never stretched."""
    p = FIGS / (name + ".png")
    w_px, h_px = Image.open(p).size
    ar = w_px / h_px
    w, h = max_w, max_w / ar
    if h > max_h:
        h, w = max_h, max_h * ar
    if centre:
        x = x + (max_w - w) / 2
    slide.shapes.add_picture(str(p), Inches(x), Inches(y + (max_h - h) / 2),
                             width=Inches(w), height=Inches(h))


def line(text):
    """A strip body: one line, ink."""
    return [[(text, {"color": INK})]]


# ============================================================ the table shape
#
# WHY THIS IS WRITTEN HERE AND NOT IN THE FURNITURE. The furniture carries no table
# helper because the deck it was cut for has no tables on it. This one has seven, so
# the shape is written once and every table on the deck is this call.
#
# IT RAISES RATHER THAN OVERFLOWS, on the same principle as check_fits and
# titled_panel. A cell that is a tenth of an inch too wide prints across its
# neighbour's column and reads as a table that lost a value, and nobody catches it
# until the export. Measured against Calibri's own advance widths, not counted.
#
# THE COLOUR VOCABULARY IS THREE VALUES AND NO MORE. Header text is MUTED over an
# INK rule. Data rules are LINE. A row or a cell that the page wants found is filled
# TINT, which is the same signal the highlight strip uses, so a tinted row means
# "this is the row the headline is about" and nothing else. ACCENT appears on a
# verdict chip and nowhere inside a table.

CELL_PAD = 0.05                 # left and right inset inside a cell, each side
RULE_H = 0.008                  # a data rule
HEAD_RULE_H = 0.013             # the heavier rule under the header row


def _cell(slide, x, y, w, h, text, pt, bold, colour, right, wrap):
    tb = box(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = MSO_ANCHOR.TOP if wrap else MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(CELL_PAD)
    tf.margin_right = Inches(CELL_PAD)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    set_text(tf, [[(text, {"size": pt, "bold": bold, "color": colour})]],
             align=PP_ALIGN.RIGHT if right else PP_ALIGN.LEFT,
             space_after=0, line_spacing=1.06)
    return tb


def _rule(slide, x, y, w, h, colour):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = colour
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _fill(slide, x, y, w, h, colour):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = colour
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def table(slide, x, y, widths, header, rows, row_h=0.28, head_h=0.28, align=None,
          pt=9.5, head_pt=9.0, row_fill=None, cell_fill=None, bold_rows=(),
          rail=None, rail_w=0.0, wrap=False, where="table"):
    """A header row and data rows laid out from panel primitives.

    widths are inches per column and must be given for every column. align is one
    character per column, l or r, and numeric columns take r. row_fill maps a row
    index to a fill colour, cell_fill maps (row, column) to one. rail is a list of
    (label, row count) that names groups of rows in a column to the left of the
    table, which is how a table gets a grouping without spending a row on each
    group name.

    Returns the y of the table's bottom edge.
    """
    ncol = len(widths)
    if len(header) != ncol:
        raise ValueError("{}: {} headers against {} columns.".format(
            where, len(header), ncol))
    align = align or ("l" + "r" * (ncol - 1))
    if len(align) != ncol:
        raise ValueError("{}: alignment string is {} against {} columns.".format(
            where, len(align), ncol))
    row_fill = row_fill or {}
    cell_fill = cell_fill or {}
    dx = x + rail_w
    total = sum(widths)
    inner = [w - 2 * CELL_PAD for w in widths]

    # MEASURE EVERYTHING FIRST. A table that is going to fail should fail before it
    # has put half of itself on the page.
    for j, h in enumerate(header):
        check_fits(h, head_pt, inner[j], "{} header {!r}".format(where, h), True)
    for i, r in enumerate(rows):
        if len(r) != ncol:
            raise ValueError("{}: row {} has {} cells against {} columns.".format(
                where, i, len(r), ncol))
        for j, c in enumerate(r):
            if wrap:
                n = wrapped_lines(str(c), pt, inner[j])
                need = 0.04 + n * pt * LINE_PITCH_IN
                if need > row_h:
                    raise ValueError(
                        "{}: cell ({}, {}) wraps to {} lines and stacks {:.3f} in "
                        "against the {:.2f} in row height, so it would print on the "
                        "row below. Shorten it or give the rows more height. Text: "
                        "{!r}".format(where, i, j, n, need, row_h, str(c)))
            else:
                check_fits(str(c), pt, inner[j],
                           "{} cell ({}, {})".format(where, i, j))
    if rail is not None:
        if sum(n for _, n in rail) != len(rows):
            raise ValueError("{}: the rail covers {} rows against {}.".format(
                where, sum(n for _, n in rail), len(rows)))

    # HEADER. Text over a rule, with no fill, so the tint stays available to mean
    # one thing only: the row this page is about.
    cx = dx
    for j, h in enumerate(header):
        _cell(slide, cx, y, widths[j], head_h, h, head_pt, True, MUTED,
              align[j] == "r", False)
        cx += widths[j]
    _rule(slide, x, y + head_h, total + rail_w, HEAD_RULE_H, INK)

    ry = y + head_h + HEAD_RULE_H
    for i, r in enumerate(rows):
        if i in row_fill:
            _fill(slide, dx, ry, total, row_h, row_fill[i])
        cx = dx
        for j, c in enumerate(r):
            if (i, j) in cell_fill:
                _fill(slide, cx, ry, widths[j], row_h, cell_fill[(i, j)])
            _cell(slide, cx, ry, widths[j], row_h, str(c), pt, i in bold_rows,
                  INK, align[j] == "r", wrap)
            cx += widths[j]
        _rule(slide, dx, ry + row_h, total, RULE_H, LINE)
        # THE SEPARATOR GOES WHITE WHERE THE FILL IS. The house rule colour is
        # grey_lighter, which disappears against the tint, so two tinted rows one
        # above the other read as a single block and the reader loses a row boundary
        # exactly where the page is asking them to look.
        if i in row_fill:
            _rule(slide, dx, ry + row_h, total, RULE_H, WHITE)
        else:
            cx = dx
            for j in range(ncol):
                if (i, j) in cell_fill:
                    _rule(slide, cx, ry + row_h, widths[j], RULE_H, WHITE)
                cx += widths[j]
        ry += row_h + RULE_H

    # THE RAIL, drawn last so its group rules sit over the data rules they replace.
    if rail is not None:
        gy = y + head_h + HEAD_RULE_H
        for k, (lab, n) in enumerate(rail):
            gh = n * (row_h + RULE_H) - RULE_H
            lines = wrapped_lines(lab, pt, rail_w - 2 * CELL_PAD, True)
            tb = box(slide, x, gy, rail_w - CELL_PAD, gh)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.0)
            tf.margin_right = Inches(CELL_PAD)
            set_text(tf, [[(lab, {"size": pt, "bold": True, "color": MUTED})]],
                     space_after=0, line_spacing=1.06)
            if lines * pt * LINE_PITCH_IN > gh:
                raise ValueError("{}: rail label {!r} does not fit its group.".format(
                    where, lab))
            if k:
                _rule(slide, x, gy - RULE_H, total + rail_w, RULE_H, MUTED)
            gy += n * (row_h + RULE_H)
    return ry


# ---------------------------------------------------------- the smaller shapes


def block_head(slide, x, y, w, text, pt=11.5):
    """The name over a statistics block or one half of a two-table page."""
    check_fits(text, pt, w, "block head")
    tb = box(slide, x, y, w, 0.26)
    set_text(tb.text_frame, [[(text, {"size": pt, "bold": True, "color": INK})]],
             space_after=0, line_spacing=1.06)
    return tb


def verdict(slide, x, y, w, text, pt=11.5):
    """A tinted chip carrying the verdict on the hypothesis the page tests."""
    check_fits(text, pt, w - 0.34, "verdict chip", True)
    band = _fill(slide, x, y, w, 0.30, TINT)
    _fill(slide, x, y, 0.055, 0.30, ACCENT)
    tb = box(slide, x + 0.20, y + 0.025, w - 0.24, 0.26)
    set_text(tb.text_frame, [[(text, {"size": pt, "bold": True, "color": INK})]],
             space_after=0, line_spacing=1.06)
    return band


def stack_panel(slide, x, y, w, h, paras, gap_pt=3, where="panel"):
    """A bordered panel carrying several measured paragraphs.

    paras is a list of (text, point size, bold, colour). The furniture's
    titled_panel takes exactly a lead-in and one body paragraph at 13 pt, and half
    the pages here want two paragraphs or a smaller size, so this is the general
    case of it and it raises on the same principle.
    """
    col = w - PANEL_ML - PANEL_MR
    spec = []
    for k, (text, pt, bold, _c) in enumerate(paras):
        spec.append((wrapped_lines(text, pt, col, bold), pt,
                     0 if k == len(paras) - 1 else gap_pt))
    need = stack_in(spec, 0.13, PANEL_MB)
    if need > h:
        raise ValueError(
            "{}: the text stacks {:.3f} in against the {:.2f} in the box has, so its "
            "last line would print on the bottom rule. Shorten it or give the box "
            "more height. Text: {!r}".format(where, need, h, paras[0][0][:52]))
    sh = panel(slide, x, y, w, h)
    set_text(sh.text_frame,
             [[(t, {"size": p, "bold": b, "color": c})] for t, p, b, c in paras],
             space_after=0, line_spacing=1.06)
    for k in range(len(paras) - 1):
        sh.text_frame.paragraphs[k].space_after = Pt(gap_pt)
    return sh


def bare_lines(slide, x, y, w, h, paras, gap_pt=3, where="lines"):
    """The same stack with no border, for a note that sits under a table."""
    spec = [(wrapped_lines(t, p, w, b), p, 0 if k == len(paras) - 1 else gap_pt)
            for k, (t, p, b, _c) in enumerate(paras)]
    need = stack_in(spec, 0.02, 0.02)
    if need > h:
        raise ValueError("{}: the text stacks {:.3f} in against {:.2f} in. "
                         "Text: {!r}".format(where, need, h, paras[0][0][:52]))
    tb = box(slide, x, y, w, h)
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame,
             [[(t, {"size": p, "bold": b, "color": c})] for t, p, b, c in paras],
             space_after=0, line_spacing=1.06)
    for k in range(len(paras) - 1):
        tb.text_frame.paragraphs[k].space_after = Pt(gap_pt)
    return tb


def stat_band(slide, y, h, cells, num_pt=19.0, lab_pt=10.5):
    """A row of bordered cells, each a figure over what the figure counts."""
    n = len(cells)
    w = (12.88 - 0.43 - BAND_GAP * (n - 1)) / n
    for i, (num, lab) in enumerate(cells):
        x = 0.43 + i * (w + BAND_GAP)
        col = w - PANEL_ML - PANEL_MR
        need = stack_in([(1, num_pt, 3),
                         (wrapped_lines(lab, lab_pt, col), lab_pt, 0)], 0.10, 0.08)
        if need > h:
            raise ValueError("statistics band cell {}: stacks {:.3f} in against "
                             "{:.2f} in. Text: {!r}".format(i, need, h, lab))
        check_fits(num, num_pt, col, "statistics band figure", True)
        sh = panel(slide, x, y, w, h, margin_top=0.10)
        sh.text_frame.margin_bottom = Inches(0.08)
        set_text(sh.text_frame,
                 [[(num, {"size": num_pt, "bold": True, "color": ACCENT})],
                  [(lab, {"size": lab_pt, "color": MUTED})]],
                 space_after=0, line_spacing=1.06)
        sh.text_frame.paragraphs[0].space_after = Pt(3)


def stat_row(slide, y, h, cells, name_pt=11.5, body_pt=12.0):
    """A row of bordered cells, each a measure over what happened to it."""
    n = len(cells)
    w = (12.88 - 0.43 - BAND_GAP * (n - 1)) / n
    for i, (name, body) in enumerate(cells):
        stack_panel(slide, 0.43 + i * (w + BAND_GAP), y, w, h,
                    [(name, name_pt, True, MUTED), (body, body_pt, False, INK)],
                    where="statistics row cell {}".format(i))


def page(title, subtitle, headline, strip, note=None, src=SRC):
    """Open a page and register it. The strip and the chrome close it."""
    s = plain_slide(title, subtitle)
    chrome(s)
    return s


def close(s, title, headline, strip, note=None, src=SRC):
    highlight_strip(s, headline, line(strip))
    if note:
        footnote(s, note)
    source(s, src)
    RUNNING_ORDER.append((title, headline))


# ===================================================================== 1. why
T = "What rides on a district's crop mix"
s = page(T, "What a district's crop mix decides, and what this settles", "", "")
stack_panel(s, STRIP_X, 1.38, STRIP_W, 0.94, [(
    "A district that puts four fifths of its cropped area under one crop is exposed "
    "to that crop's price and that crop's water demand. The mix sits underneath "
    "groundwater draw, procurement exposure and what a farming household earns in a "
    "bad year. District statistics report area by crop every year and carry no "
    "measure of how that area is split.", 13.0, False, INK)], where="s1 lede")
titled_panel(s, THIRD_X[0], 2.38, THIRD_W, 1.96, "What was open",
             "Whether a district's crop mix follows the water it has, the markets it "
             "can reach or its agrarian structure. Cropping statistics carry area by "
             "crop and village records carry the rest, and the two have not been read "
             "together.", where="s1 a")
titled_panel(s, THIRD_X[1], 2.38, THIRD_W, 1.96, "The record",
             "Ministry of Agriculture district statistics: area and production by crop "
             "and season, 54 crops, 23 agricultural years, 14,089 district-year "
             "records. Village records on irrigation, markets, input supply and "
             "agrarian structure, aggregated to district.", where="s1 b")
titled_panel(s, THIRD_X[2], 2.38, THIRD_W, 1.96, "What this settles",
             "A measure of the crop mix that reads as a count of crops, the national "
             "picture and its geography, thirteen candidate explanations screened on "
             "one sample, and five taken forward with the coefficient that decides "
             "each.", where="s1 c")
stat_band(s, 4.46, 0.78, [("725", "districts with a cropping record"),
                          ("606", "carrying village records too"),
                          ("54", "crops"),
                          ("23", "agricultural years")])
close(s, T, "Half of a district's cropped area sits under its single largest crop",
      "The largest crop takes 50.6 percent of cropped area across the 725 districts, "
      "and the top three take 80.3 percent.",
      "The join matches 99.6 percent of the cropping record, and 64 of the 606 "
      "districts sit on a shared pre-2011 parent, which every estimate is tested "
      "against.")

# ============================================================== 2. resolution
T = "Counting crops and weighting them"
s = page(T, "The effective number of crops against crops grown in an average year, "
            "by district", "", "")
place(s, "f1_count_vs_effective", TWOUP_X[0], 1.50, TWOUP_W, 3.40)
table(s, TWOUP_X[1], 1.42, [2.30, 0.65, 3.105],
      ["Measure", "Mean", "What it reads as"],
      [["Crops grown in an average year", "21.2", "crops with any recorded area"],
       ["Effective number of crops", "4.9",
        "equally common crops giving the same diversity"],
       ["Effective number weighted to the dominant crops", "3.5",
        "the same, with rare crops discounted further"],
       ["Evenness", "0.254", "the effective number over the count"]],
      row_h=0.38, head_h=0.28, align="lrl", pt=9.0, wrap=True, where="s2 scale")
stack_panel(s, TWOUP_X[1], 3.30, TWOUP_W, 1.95, [
    ("Across the 606 districts the effective number of crops runs 2.4 at the tenth "
     "percentile, 4.3 at the median and 8.1 at the ninetieth.", 13.0, False, INK),
    ("The effective number of crops is the exponential of the Shannon entropy of a "
     "district's cropped-area shares, taken within each year and then averaged. It "
     "reads as a count. A district growing twenty crops with four fifths of its land "
     "under paddy farms like one growing three, and the measure says three.",
     13.0, False, INK)], where="s2 note")
close(s, T, "The average district is effectively growing five crops",
      "It grows 21.2 crops in an average year, and evenness across the 606 districts "
      "averages 0.254.",
      "Measured within an agricultural year and then averaged. The composite index "
      "used in earlier work is a min-max blend whose district scores move by up to "
      "0.0156 when the sample changes.")

# ============================================================== 3. resolution
T = "The national picture in numbers"
s = page(T, "Cropping, composition and rural structure across 606 districts in 30 "
            "states", "", "")
AX, BX, CX = 0.43, 5.93, 9.33
AW, BW, CW = 5.30, 3.20, 3.55
block_head(s, AX, 1.36, AW, "Cropping, across the 606 districts")
table(s, AX, 1.66, [2.30, 0.75, 0.75, 0.75, 0.75],
      ["Measure", "Mean", "10th", "Median", "90th"],
      [["Crops grown in an average year", "21.2", "11.1", "22.4", "29.3"],
       ["Effective number of crops", "4.9", "2.4", "4.3", "8.1"],
       ["Evenness", "0.254", "0.137", "0.242", "0.382"],
       ["Irrigation share of sown area", "0.546", "0.235", "0.529", "0.879"],
       ["Cropping intensity", "1.75", "1.44", "1.77", "2.03"],
       ["Mean holding size, hectares", "1.96", "0.61", "1.35", "2.95"]],
      row_h=0.30, where="s3 block a")
stack_panel(s, AX, 3.86, AW, 1.39, [
    ("Eight districts in ten sit between 0.24 and 0.88 on the irrigation share of "
     "sown area, so the sample spans the range of Indian water availability rather "
     "than one end of it. Mean holding size is inferred from village records rather "
     "than measured. It is carried through every estimate on this deck and is not "
     "distinguishable from zero in any of them.",
     12.0, False, INK)], where="s3 block a note")
block_head(s, BX, 1.36, BW, "Share of national cropped area")
table(s, BX, 1.66, [1.95, 1.25], ["Category", "Share"],
      [["cereals", "56.8"], ["oilseeds", "16.2"], ["pulses", "12.9"],
       ["fibre crops", "7.0"], ["sugar", "2.5"], ["fodder", "1.6"],
       ["spices", "1.2"], ["vegetables", "1.1"], ["fruit", "0.4"]],
      row_h=0.30, row_fill={0: TINT}, where="s3 block b")
block_head(s, CX, 1.36, CW, "Share of a district's villages")
table(s, CX, 1.66, [2.30, 1.25], ["Facility", "Mean share"],
      [["producer organisation", "0.231"], ["fertiliser shop", "0.175"],
       ["weekly haat", "0.143"], ["farm-gate processing", "0.114"],
       ["regular market", "0.098"], ["seed centre", "0.098"],
       ["custom hiring", "0.092"], ["cold storage", "0.084"],
       ["soil testing", "0.043"], ["mandi", "0.030"]],
      row_h=0.29, where="s3 block c")
bare_lines(s, CX, 5.01, CW, 0.28,
           [("40 districts have no mandi and 30 have no weekly haat.",
             10.5, False, MUTED)], where="s3 block c tail")
close(s, T, "Cereals take 57 percent of India's cropped area and fruit and "
            "vegetables take 1.5 percent",
      "A mandi sits in 3 of every hundred villages and a weekly haat in 14.",
      "Percentiles are across the 606 districts. Crop shares are weighted by each "
      "district's mean annual cropped area. Village shares come from the 2019 "
      "village records.")

# ============================================================== 4. resolution
T = "Effective number of crops by district"
s = page(T, "Averaged over each district's years, 1997-98 to 2019-20", "", "")
place(s, "f12_map_effective_large", 0.43, 1.24, 6.05, 4.02)
RAIL_X, RAIL_W2 = 6.90, 5.98
# The two state tables sit side by side inside the rail beside the map, so the
# district count comes out and the note below carries it instead.
TW = 2.80
block_head(s, RAIL_X, 1.34, TW, "Widest")
table(s, RAIL_X, 1.62, [1.72, 1.08],
      ["State", "Effective crops"],
      [["Karnataka", "8.0"], ["Andhra Pradesh", "6.9"], ["Nagaland", "6.4"],
       ["Uttarakhand", "6.2"], ["Madhya Pradesh", "6.2"]],
      row_h=0.30, where="s4 widest")
block_head(s, RAIL_X + TW + 0.38, 1.34, TW, "Narrowest")
table(s, RAIL_X + TW + 0.38, 1.62, [1.72, 1.08],
      ["State", "Effective crops"],
      [["Tripura", "2.0"], ["Puducherry", "2.4"], ["Odisha", "2.5"],
       ["Punjab", "2.8"], ["Mizoram", "2.8"]],
      row_h=0.30, where="s4 narrowest")
stack_panel(s, RAIL_X, 3.52, RAIL_W2, 1.73, [
    ("Chhattisgarh grows 31.6 crops in an average year, more than any other state, "
     "and is effectively growing 3.2 on an evenness of 0.107.", 12.0, False, INK),
    ("Punjab reaches the narrow end through wheat and Odisha through paddy. Both "
     "grow enough crops, and the land sits under one of them.", 12.0, False, INK),
    ("The dark districts run through the peninsular and central belt, and the pale "
     "ones through the eastern rice belt and the irrigated northwest.",
     12.0, False, INK)], where="s4 rail note")
close(s, T, "Karnataka districts average 8.0 effective crops, the highest of any state",
      "Odisha districts average 2.5 and Punjab districts 2.8, reaching the same "
      "place through paddy and through wheat.",
      "State figures are unweighted means of their districts, over 30 districts in "
      "Karnataka, 55 in Madhya Pradesh, 30 in Odisha and 22 in Punjab. Grey "
      "districts carry no agricultural record.")

# ============================================================== 5. resolution
T = "The crop that takes the most land"
s = page(T, "Largest single land use by district, and what each quarter of districts "
            "plants", "", "")
block_head(s, STRIP_X, 1.30, 6.10, "Largest land use, of the 725 districts")
HEAD5 = ["Crop", "Districts", "Share of cropped area"]
table(s, 0.43, 1.56, [2.60, 1.10, 2.40], HEAD5,
      [["Rice", "350", "59.5%"], ["Wheat", "147", "42.1%"],
       ["Maize", "36", "45.7%"], ["Cotton", "33", "38.7%"],
       ["Soyabean", "28", "39.5%"], ["Coconut", "19", "64.8%"]],
      row_h=0.25, head_h=0.26, row_fill={0: TINT}, where="s5 crops left")
table(s, 6.78, 1.56, [2.60, 1.10, 2.40], HEAD5,
      [["Jowar", "19", "28.5%"], ["Bajra", "16", "40.4%"],
       ["Groundnut", "12", "47.0%"], ["Rapeseed and mustard", "10", "36.5%"],
       ["Sugarcane", "10", "45.2%"], ["Ragi", "9", "39.0%"]],
      row_h=0.25, head_h=0.26, where="s5 crops right")
block_head(s, STRIP_X, 3.46, 8.00,
           "Composition by quarter of the effective number of crops, percent of "
           "cropped area")
table(s, STRIP_X, 3.76, [2.45] + [1.25] * 8,
      ["Quarter", "cereals", "oilseeds", "pulses", "fibre crops", "spices", "sugar",
       "vegetables", "fruit"],
      [["Least diverse quarter", "76.7", "10.3", "5.2", "2.0", "3.7", "0.9", "0.8",
        "0.3"],
       ["Second", "70.0", "11.2", "5.6", "4.5", "1.8", "3.4", "1.7", "1.2"],
       ["Third", "63.6", "13.8", "10.4", "4.1", "2.5", "2.1", "2.2", "0.5"],
       ["Most diverse quarter", "49.9", "16.6", "17.3", "7.5", "2.3", "2.3", "2.0",
        "1.1"]],
      row_h=0.29, cell_fill={(i, 1): TINT for i in range(4)}, where="s5 quarters")
close(s, T, "Rice or wheat is the largest land use in 497 of 725 districts",
      "Where cereal area gives way it goes to pulses and oilseeds, and fruit and "
      "vegetables stay near a fiftieth of the land.",
      "Districts are sorted into quarters on the effective number of crops across the "
      "725 with a cropping record. Category shares are averaged within the quarter.")

# ============================================================== 6. resolution
T = "Crop diversity over two decades"
s = page(T, "Fixed panel of 429 districts reporting in every year, 1998 to 2019",
         "", "")
place(s, "f13_trend_wide", STRIP_X, 1.40, STRIP_W, 2.90)
stat_row(s, 4.42, 0.80, [
    ("Effective number of crops", "165 districts up, 264 down, median -0.19"),
    ("Crops grown in an average year", "199 up, 226 down, median -0.29"),
    ("Evenness", "200 up, 229 down, median -0.004")])
close(s, T, "The effective number of crops held at about 5.4 across two decades",
      "Crops grown in an average year rose from 21.2 to 24.9 and evenness fell from "
      "0.265 to 0.224, so the added crops sit on small patches.",
      "The panel holds the same 429 districts in every year, unweighted across "
      "districts. Each district's 2013 to 2019 mean is compared against its 1998 to "
      "2004 mean.")

# ================================================== 7. identifying hypotheses
T = "Crop diversity against irrigation"
s = page(T, "Effective number of crops by irrigation share of sown area, 606 "
            "districts", "", "")
place(s, "f5_irrigation_shape", TWOUP_X[0], 1.55, TWOUP_W, 3.30)
table(s, TWOUP_X[1], 1.42, [0.75, 0.85, 0.98, 1.20, 1.15, 1.125],
      ["Decile", "Districts", "Irrigation", "Effective crops", "Crops grown",
       "Cereal share"],
      [["1", "61", "0.165", "4.11", "19.1", "0.718"],
       ["2", "61", "0.300", "4.99", "21.2", "0.669"],
       ["3", "60", "0.385", "5.29", "21.9", "0.617"],
       ["4", "61", "0.444", "5.86", "21.1", "0.638"],
       ["5", "60", "0.502", "5.01", "20.5", "0.618"],
       ["6", "61", "0.563", "5.63", "22.1", "0.587"],
       ["7", "60", "0.640", "5.24", "21.5", "0.574"],
       ["8", "61", "0.718", "4.96", "24.3", "0.604"],
       ["9", "60", "0.826", "4.54", "22.0", "0.772"],
       ["10", "61", "0.922", "3.71", "18.0", "0.793"]],
      row_h=0.29, align="rrrrrr", row_fill={3: TINT, 9: TINT}, where="s7 deciles")
bare_lines(s, TWOUP_X[1], 4.75, TWOUP_W, 0.52, [
    ("Deciles are of the irrigation share of sown area across the 606 districts, "
     "1 = driest and 10 = wettest, about 61 districts each.", 11.5, False, MUTED)],
    where="s7 decile note")
close(s, T, "The effective number of crops rises with the irrigation share of sown "
            "area and then falls",
      "Districts in the wettest tenth average 3.71 effective crops against 5.86 in "
      "the fourth tenth.",
      "The fitted curve carries no state fixed effects and turns at 51 percent of "
      "sown area. State fixed effects move it to 28 percent, and across the eight "
      "samples it runs 25 to 51 percent.")

# ================================================== 8. identifying hypotheses
T = "Thirteen candidate explanations"
s = page(T, "Association with the effective number of crops, before and after "
            "adjustment. One star is p below 0.10, two below 0.05, three below 0.01",
         "", "")
table(s, STRIP_X, 1.42, [2.95, 0.75, 1.05, 1.00, 5.10],
      ["Candidate", "Districts", "Unadjusted", "Adjusted", "Where it comes from"],
      [["irrigation share of sown area", "606", "-0.753**", "+3.901***",
        "the shape on the previous page"],
       ["canal village share", "606", "-1.265***", "-0.246",
        "command-area rotations"],
       ["mean holding size, log hectares", "606", "+0.339***", "+0.019",
        "small farms spreading risk"],
       ["cultivator share of agricultural workers", "606", "-1.134***", "-0.060",
        "labour supply and crop choice"],
       ["landless share", "592", "-0.204", "-0.939", "rural deprivation records"],
       ["mandi village share", "606", "-1.931", "-0.142", "assured procurement"],
       ["weekly haat village share", "606", "-0.148", "+0.857",
        "smallholder selling points"],
       ["regular market village share", "606", "-2.307***", "-0.716",
        "daily retail demand"],
       ["producer organisation village share", "606", "+2.188***", "+1.905**",
        "collective aggregation"],
       ["cold storage village share", "606", "+6.308***", "+1.335",
        "perishables need cold"],
       ["SC population share", "606", "+1.368", "+1.233", "caste and land access"],
       ["ST population share", "606", "-0.350", "+0.730", "forest-fringe cropping"],
       ["cropping intensity", "606", "+0.845**", "+2.265***",
        "an accounting identity as much as a lever"]],
      row_h=0.265, align="lrrrl", pt=9.0, head_pt=9.0, bold_rows=(0, 8, 12),
      rail=[("Water", 2), ("Agrarian structure", 3),
            ("Markets and infrastructure", 5), ("Population and land use", 3)],
      rail_w=1.60, where="s8 screen")
close(s, T, "Three of the thirteen candidates keep a significant coefficient after adjustment",
      "The irrigation share of sown area, producer organisation coverage and "
      "cropping intensity.",
      "Adjusted models add state fixed effects, log mean holding size and the "
      "irrigation share of sown area with its square. Cropping intensity matches "
      "published state figures at 0.65.")

# ========================================================= 9. the hypotheses
T = "The five hypotheses and what would falsify each"
s = page(T, "Set down before the estimates, with the test each one gets", "", "")
table(s, STRIP_X, 1.44, [1.75, 3.75, 3.60, 3.35],
      ["Hypothesis", "Mechanism", "What would confirm it", "What would falsify it"],
      [["H1 Water quantity",
        "Irrigation lifts the constraint on what can be grown at low levels, and at "
        "high levels it pays for one assured package",
        "A negative squared term on the irrigation share of sown area that holds "
        "across sample and index choice",
        "A null or positive squared term, or a sign that depends on which index is "
        "used"],
       ["H2 Water type",
        "Canal water arrives on a rotation the system sets and groundwater arrives "
        "when the cultivator chooses, so fewer crops fit the canal calendar",
        "Source coefficients that stay large with the level of irrigation held "
        "constant",
        "Null source coefficients once the level of irrigation is in"],
       ["H3 Somewhere to sell",
        "A weekly haat lets a smallholder sell small volumes of many things, so crops "
        "with no other outlet get planted",
        "Haat coverage raising crops grown in an average year, surviving development "
        "controls",
        "Null haat coefficients, or coefficients that vanish under those controls"],
       ["H4 Collective marketing",
        "A producer organisation aggregates volume for crops a single farm cannot "
        "market, so land shifts between crops already grown",
        "Producer organisation coverage raising the effective number of crops with "
        "crops grown flat",
        "Both institutions moving both measures together"],
       ["H5 Assured procurement",
        "Dense procurement and input supply make the cereal package the only paying "
        "choice once water is plentiful, which is what pulls the curve down after "
        "the peak",
        "A steeper fall after the peak where mandi and input coverage is dense, and a "
        "higher cereal share where mandis are dense",
        "A positive interaction, or a negative mandi coefficient on cereal share"]],
      row_h=0.66, head_h=0.30, align="llll", pt=9.5, wrap=True, where="s9 hypotheses")
close(s, T, "Each hypothesis names the coefficient that decides it",
      "The outcome is the effective number of crops throughout, with state fixed "
      "effects on every estimate.",
      "Adjusted models carry the irrigation share of sown area and its square, log "
      "mean holding size, ST population share, log night lights, non-farm "
      "establishment density and a connectivity index.")

# =============================================================== 10. findings
T = "Fourteen versions of the irrigation estimate"
s = page(T, "H1, eight variations of the specification and six alternative indices",
         "", "")
verdict(s, STRIP_X, 1.32, 1.70, "H1 holds")
table(s, STRIP_X, 1.64, [3.90, 1.05, 1.40, 1.45, 1.45, 1.10, 1.00],
      ["Specification", "Districts", "Linear term", "Squared term", "p on squared",
       "Turn", "R-squared"],
      [["No fixed effects", "606", "+10.444", "-10.185", "<0.001", "0.513", "0.065"],
       ["State fixed effects", "606", "+3.874", "-6.877", "<0.001", "0.282", "0.530"],
       ["Irrigation from the 2011 village directory", "573", "+2.749", "-5.428",
        "<0.001", "0.253", "0.534"],
       ["Drop shared pre-2011 parents", "542", "+3.396", "-6.382", "<0.001", "0.266",
        "0.538"],
       ["Inverse-parent weighted", "606", "+3.641", "-6.612", "<0.001", "0.275",
        "0.531"],
       ["Drop source-conflict districts", "530", "+3.235", "-6.615", "<0.001",
        "0.245", "0.547"],
       ["Full 23 years only", "506", "+3.107", "-6.270", "<0.001", "0.248", "0.516"],
       ["Weighted by cropped area", "606", "+5.778", "-8.083", "<0.001", "0.357",
        "0.533"],
       ["Crops grown in an average year", "606", "+10.088", "-10.910", "0.0001",
        "0.462", "0.828"],
       ["Effective number weighted to the dominant crops", "606", "+1.760", "-3.925",
        "0.0001", "0.224", "0.500"],
       ["Evenness", "606", "+0.066", "-0.188", "0.013", "0.177", "0.543"],
       ["Shannon", "606", "+0.936", "-1.432", "<0.001", "0.327", "0.578"],
       ["Simpson", "606", "+0.268", "-0.367", "0.0003", "0.364", "0.611"],
       ["Composite index, for comparison", "606", "+0.310", "-0.417", "<0.001",
        "0.372", "0.668"]],
      row_h=0.232, head_h=0.27, pt=9.0, head_pt=9.0,
      rail=[("specification", 8), ("index", 6)], rail_w=1.10, where="s10 robustness")
close(s, T, "The squared term on the irrigation share of sown area is negative in "
            "all fourteen estimates",
      "Eight of them vary the specification and six substitute a different index, "
      "and the weakest still carries p = 0.013.",
      "The turn is the maximum of the fitted curve, on the irrigation share of sown "
      "area. Crops grown pooled across years shows no hump at p = 0.19 and within "
      "each year it does at p = 0.0001.")

# =============================================================== 11. findings
T = "Irrigation source and what a district grows"
s = page(T, "H2, canal and surface water against groundwater at the same irrigation "
            "share", "", "")
place(s, "f7_irrigation_source", TWOUP_X[0], 1.60, TWOUP_W, 3.20)
verdict(s, TWOUP_X[1], 1.36, 2.10, "H2 holds")
table(s, TWOUP_X[1], 1.74, [2.85, 0.88, 0.62, 1.02, 0.685],
      ["Outcome", "Canal", "p", "Surface water", "p"],
      [["Effective number of crops", "-0.310", "0.42", "-2.030", "0.026"],
       ["Crops grown in an average year", "-2.593", "0.003", "-0.467", "0.73"],
       ["Effective number weighted to the dominant crops", "-0.155", "0.59",
        "-1.660", "0.019"],
       ["Evenness", "+0.034", "0.073", "-0.087", "0.11"],
       ["Cereal share", "+0.022", "0.60", "+0.300", "0.0001"],
       ["Pulse share", "+0.073", "0.0005", "-0.104", "0.002"],
       ["Oilseed share", "-0.097", "0.002", "-0.164", "0.008"]],
      row_h=0.29, pt=9.0, where="s11 source")
stack_panel(s, TWOUP_X[1], 4.11, TWOUP_W, 1.14, [
    ("The irrigation share of sown area and its square are held in every row, so "
     "these are differences between water sources at the same amount of water.",
     12.0, False, INK),
    ("Groundwater is the omitted source and covers 50.5 percent of villages, other "
     "sources 21.1, canal 17.2 and surface water 11.3.", 12.0, False, INK)],
    where="s11 note")
close(s, T, "Surface-water districts average two fewer effective crops at the same "
            "irrigation share",
      "Canal dependence takes 2.6 crops off the list, adds 0.073 to the pulse "
      "share and leaves the effective number alone.",
      "Coefficients are on the share of a district's villages naming that source, so "
      "they read from none to all. State fixed effects, errors clustered on the "
      "pre-2011 parent district.")

# =============================================================== 12. findings
T = "Rural facilities and cropping"
s = page(T, "H3 and H4, each facility entered on its own with the full control set",
         "", "")
place(s, "f10_market_solo", TWOUP_X[0], 1.55, TWOUP_W, 2.70)
stack_panel(s, TWOUP_X[0], 4.35, TWOUP_W, 0.90, [
    ("A one standard deviation rise in haat coverage, 0.117 of a district's "
     "villages, goes with about 0.9 more crops against a mean of 21.2, and the "
     "coefficient runs +6.44 to +8.02 across every control specification tried.",
     12.0, False, INK)], where="s12 scale")
verdict(s, TWOUP_X[1], 1.36, 4.10, "H3 holds on crops grown only; H4 holds")
table(s, TWOUP_X[1], 1.74, [1.95, 1.15, 1.10, 0.95, 0.905],
      ["Facility", "Effective crops", "Crops grown", "Cereal share", "Pulse share"],
      [["Weekly haat", "+0.49 (0.56)", "+7.33 (0.000)", "-0.050 (0.42)", "-0.068 (0.11)"],
       ["Cold storage", "+0.75 (0.65)", "+6.34 (0.039)", "+0.141 (0.13)", "-0.172 (0.001)"],
       ["Seed centre", "-0.38 (0.71)", "+4.53 (0.022)", "-0.010", "-0.057"],
       ["Farm-gate processing", "+0.20 (0.82)", "+3.96 (0.047)", "+0.035", "-0.053"],
       ["Soil testing", "+0.32 (0.75)", "+2.67 (0.17)", "+0.111", "-0.075"],
       ["Producer organisation", "+1.95 (0.032)", "+2.39 (0.19)", "-0.134*",
        "-0.066*"],
       ["Mandi", "+0.23 (0.91)", "+1.88 (0.54)", "-0.235*", "-0.046"],
       ["Custom hiring", "+0.11 (0.93)", "+1.75 (0.48)", "-0.032", "-0.056"],
       ["Fertiliser shop", "+0.28 (0.78)", "+1.10 (0.57)", "-0.198 (0.021)", "-0.174 (0.000)"],
       ["Regular market", "-0.69 (0.39)", "-2.88 (0.12)", "+0.049 (0.55)", "-0.094 (0.004)"]],
      row_h=0.29, pt=9.0, row_fill={0: TINT, 5: TINT}, where="s12 solo")
close(s, T, "Weekly haat coverage goes with 7.3 more crops grown in an average year",
      "Producer organisation coverage adds 1.95 to the effective number of crops and "
      "leaves crops grown flat.",
      "Nine of the ten facilities correlate 0.24 to 0.83 with one another and the "
      "weekly haat runs -0.21 to +0.06 against them. Errors cluster on the pre-2011 "
      "parent district.")

# =============================================================== 13. findings
T = "Market density and the fall after the peak"
s = page(T, "H5, curvature of the irrigation relationship where infrastructure is "
            "thin and where it is dense", "", "")
place(s, "f14_interaction_wide", 0.43, 1.42, 8.20, 1.85)
stack_panel(s, 8.90, 1.42, 3.97, 1.85, [
    ("The other half of the same hypothesis", 12.5, True, INK),
    ("Mandi coverage was expected to raise the cereal share. Entered alone with the "
     "full control set it goes with a cereal share 0.235 lower (p = 0.054), and "
     "entered with the other output markets, 0.240 lower.", 12.0, False, INK)],
    where="s13 callout")
verdict(s, STRIP_X, 3.36, 2.15, "H5 does not hold")
bare_lines(s, 2.90, 3.40, 9.98, 0.30, [
    ("Districts split at the median on market density run 3.79 to 5.99 effective "
     "crops raw, and all three adjusted contrasts are null.", 10.5, False, MUTED)],
    where="s13 typology")
table(s, STRIP_X, 3.78, [3.05, 2.35, 2.40, 2.05, 2.60],
      ["Infrastructure", "Curvature where thin", "Curvature where dense",
       "Interaction", "p"],
      [["Mandi village share", "-8.98", "-3.74", "+2.62", "0.033"],
       ["Input-supply density", "-8.76", "-4.30", "+2.23", "0.056"],
       ["Post-harvest density", "-8.63", "-4.43", "+2.10", "0.087"],
       ["Output-market density", "-8.46", "-4.47", "+2.00", "0.102"],
       ["Fertiliser shop village share", "-8.20", "-5.14", "+1.53", "0.143"]],
      row_h=0.24, head_h=0.25, pt=9.5, row_fill={0: TINT}, where="s13 interaction")
close(s, T, "All five interaction terms run the opposite way to the hypothesis",
      "A positive term means the fall after the peak is gentler where "
      "infrastructure is dense, not sharper, and the mandi one is the only one "
      "under 5 percent.",
      "Curvature is the coefficient on the square of the irrigation share of sown "
      "area, read one standard deviation below and above the mean of each measure.")

# =============================================================== 14. findings
T = "The five hypotheses and their verdicts"
s = page(T, "Five hypotheses, the test applied and the number that decides each",
         "", "")
table(s, STRIP_X, 1.42, [1.85, 3.10, 5.20, 2.30],
      ["Hypothesis", "The test", "The deciding number", "Verdict"],
      [["H1 Water quantity",
        "A negative squared term on the irrigation share of sown area across "
        "fourteen specifications",
        "Negative in all fourteen, weakest p = 0.013, turn between 25 and 51 percent "
        "of sown area", "Holds"],
       ["H2 Water type",
        "Source coefficients with the irrigation share and its square held",
        "Surface water -2.03 effective crops (p = 0.026) and cereal share +0.300 "
        "(p = 0.0001); canal -2.59 crops grown (p = 0.003)", "Holds"],
       ["H3 Somewhere to sell",
        "Haat coverage on crops grown under seven specifications",
        "+6.44 to +8.02 with controls and +16.99 without, p below 0.001 throughout, "
        "with the effective number at +0.49 (p = 0.56) and evenness at -0.071",
        "Holds on crops grown only"],
       ["H4 Collective marketing",
        "Producer organisation coverage read on both measures at once",
        "+1.95 entered alone (p = 0.032) and +1.91 in the screen (p = 0.026), with "
        "crops grown flat at +2.39 (p = 0.19)", "Holds"],
       ["H5 Assured procurement",
        "A steeper fall after the peak where coverage is dense, and a higher cereal "
        "share where mandis are dense",
        "All five interactions positive, mandi +2.62 (p = 0.033), and mandi coverage "
        "on cereal share -0.235", "Does not hold"]],
      row_h=0.62, head_h=0.28, align="llll", pt=9.5, wrap=True,
      cell_fill={(0, 3): TINT, (1, 3): TINT, (2, 3): TINT, (3, 3): TINT},
      where="s14 scorecard")
bare_lines(s, STRIP_X, 4.92, STRIP_W, 0.32, [
    ("Agrarian structure was screened out before the hypotheses: mean holding size "
     "+0.02 (p = 0.88), cultivator share -0.06 (p = 0.93), landless share -0.94 "
     "(p = 0.11).", 11.0, False, MUTED)], where="s14 tail")
close(s, T, "Water quantity, water source and producer organisation coverage move "
            "the effective number of crops",
      "Mandi density, market type and agrarian structure do not.",
      "606 districts across 30 states. Irrigation systems and market facilities sit "
      "where terrain and aquifers already favoured them, and the rural record is a "
      "single 2019 cross-section.")


os.makedirs(str(OUT.parent), exist_ok=True)
prs.save(str(OUT))
print("WROTE {}".format(OUT))
print("\nRunning order, {} slides".format(len(RUNNING_ORDER)))
for i, (t, h) in enumerate(RUNNING_ORDER, 1):
    print("  {:2d}  {:46s} {}".format(i, t[:46], h))
